import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Helper function to add title slide
    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    # Helper function to add content slide
    def add_content_slide(title, content_items):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    # Helper function to add two-column slide
    def add_two_column_slide(title, left_items, right_items):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        # Left column
        p = tf.add_paragraph()
        p.text = "Left Column:"
        p.font.bold = True
        for item in left_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
        
        # Right column
        p = tf.add_paragraph()
        p.text = "\nRight Column:"
        p.font.bold = True
        for item in right_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1

    # Slide 1: Title
    add_title_slide(
        "CSCI480 Layered IDS/IPS",
        "Autonomous Intrusion Prevention System\nSenior Capstone Project (Academic Demonstration)\nJoshua Swanson | Jerry Buno | Conner Benitez"
    )

    # Slide 2: Problem Statement & Motivation
    add_content_slide("Problem Statement & Motivation", [
        "Current security systems suffer from three major weaknesses:",
        "- Passive defense: Alerts occur after damage may already be done",
        "- Static rule sets: Ineffective against zero-day or novel attacks",
        "- Alert fatigue: Excessive notifications overwhelm security teams",
        "Traditional signature-based IDS cannot detect zero-day attacks",
        "Automated attacks operate faster than human defenders can respond",
        "Organizations require layered defense with multiple detection models",
        "Need for active prevention capabilities beyond passive monitoring"
    ])

    # Slide 3: Project Objectives
    add_content_slide("Project Objectives", [
        "Learn patterns of normal network activity",
        "Detect deviations that may indicate malicious behavior using machine learning",
        "Create active response mechanisms to mitigate threats (Rule-based and AI Agent)",
        "Record and analyze incidents to support system improvement",
        "Develop an administrative dashboard for real-time monitoring and historical analysis",
        "Implement feedback loop for adaptive learning and false positive reduction"
    ])

    # Slide 4: Key Features & Functionality
    add_content_slide("Key Features & Functionality", [
        "Data Collection - Scapy packet capture with feature extraction",
        "Live packet capture with real-time analysis",
        "Layered ML detection using 7 different models (AutoEncoder, IsolationForest, KMeans, RandomForest, GradientBoosting, GNN, PPO)",
        "Active prevention and healing via firewall rules",
        "Real attack simulation for testing",
        "PCAP replay and upload for forensic analysis",
        "Real-time dashboard with visualization",
        "Adaptive decoy deployment for threat deception",
        "Feedback Loop & Adaptive Learning for continuous improvement"
    ])

    # Slide 5: System Architecture
    add_two_column_slide(
        "System Architecture",
        [
            "Backend (Flask API)",
            "- Packet capture engine (Scapy)",
            "- ML model inference pipeline",
            "- Prevention/healing logic",
            "- Attack simulator",
            "- Decoy manager"
        ],
        [
            "Frontend (Dashboard UI)",
            "- Real-time traffic visualization",
            "- Alert management interface",
            "- Defense controls",
            "- PCAP analysis tools",
            "- Results export"
        ]
    )

    # Slide 6: Technologies Used
    add_two_column_slide(
        "Technologies Used",
        [
            "Backend Frameworks",
            "- Flask (Python web framework)",
            "- Flask-CORS (cross-origin)",
            "- Scapy (packet capture)",
            "- psutil (system monitoring)"
        ],
        [
            "ML & Data Science",
            "- PyTorch (deep learning)",
            "- scikit-learn (ML models)",
            "- NumPy & Pandas (data)",
            "- 7 ML models: AE, IsoForest, KMeans, RF, GNN, PPO, GBDT"
        ]
    )

    # Slide 7: ML Models Overview
    add_content_slide("ML Models Overview", [
        "AutoEncoder (Unsupervised) - Anomaly detection via reconstruction error",
        "IsolationForest (Unsupervised) - Isolates anomalies using random partitioning",
        "KMeans (Unsupervised) - Clustering-based anomaly detection",
        "RandomForest (Supervised) - Multi-class attack classification",
        "GradientBoosting (Supervised) - Ensemble decision trees",
        "Graph Neural Network (Deep Learning) - Flow pattern analysis",
        "PPO (Reinforcement Learning) - Adaptive response policy"
    ])

    # Slide 8: Key Challenges & Solutions
    add_two_column_slide(
        "Key Challenges & Solutions",
        [
            "Challenge: Real-time packet capture",
            "Solution: Batched capture with Scapy, async processing",
            "",
            "Challenge: ML model integration",
            "Solution: Unified prediction API with feature scaling",
            "",
            "Challenge: Windows firewall integration"
        ],
        [
            "Solution: PowerShell-based firewall rule management",
            "",
            "Challenge: False positive reduction",
            "Solution: Ensemble voting with risk ranking",
            "",
            "Challenge: Deployment complexity",
            "Solution: PyInstaller packaging with Inno Setup installer"
        ]
    )

    # Slide 9: Team Roles & Lessons Learned
    add_two_column_slide(
        "Team Roles & Lessons Learned",
        [
            "Team Roles",
            "Conner Benitez",
            "- Data Pipeline & Isolation Forest",
            "- Ensures data is clean before use",
            "",
            "Joshua Swanson & Jerry Buno",
            "- Deep Learning (Autoencoder) & K-Means",
            "- Model development and integration"
        ],
        [
            "Lessons Learned",
            "Data Quality First",
            "- Spent 80% of time on Data Cleaning",
            "- Bad data breaks models instantly",
            "",
            "Model Variety",
            "- Different models do different things",
            "- Isolation Forest is fast and broad",
            "- Autoencoders are precise",
            "- Using both gives better detection"
        ]
    )

    # Slide 10: Application Screenshots
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title = slide.shapes.title
    title.text = "Application Screenshots"
    
    screenshots_dir = r"c:\Users\pompk\Desktop\CSCI480\Demo_Instructions\screenshots"
    screenshots = [
        ("screenshot_01_dashboard.png", "Dashboard Overview", "Real-time traffic monitoring interface"),
        ("screenshot_02_ui.png", "Main UI", "User interface for capture control"),
        ("screenshot_03_results.png", "Live Results", "Packet flow analysis with model predictions"),
        ("screenshot_04_defense.png", "Defense & Prevention", "Firewall rule management and decoy deployment"),
        ("screenshot_05_pcap.png", "PCAP Replay", "Forensic analysis of network traffic"),
        ("screenshot_06_models.png", "Model Settings", "ML model configuration and status")
    ]
    
    # Add screenshots in a 2x3 grid
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(2.5)
    
    for i, (filename, caption, desc) in enumerate(screenshots):
        row = i // 2
        col = i % 2
        
        img_path = os.path.join(screenshots_dir, filename)
        if os.path.exists(img_path):
            left_pos = left + (col * (width + 0.5))
            top_pos = top + (row * (height + 1.2))
            
            slide.shapes.add_picture(img_path, left_pos, top_pos, width=width, height=height)
            
            # Add caption
            textbox = slide.shapes.add_textbox(left_pos, top_pos + height + 0.1, width, Inches(0.4))
            text_frame = textbox.text_frame
            text_frame.text = caption
            text_frame.paragraphs[0].font.size = Pt(10)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Slide 11: What Makes It Unique
    add_content_slide("What Makes This Project Unique", [
        "Layered detection with 7 complementary ML models",
        "Active prevention capabilities (firewall blocking, decoy deployment)",
        "Real-time attack simulation for testing",
        "PCAP replay for forensic analysis",
        "Windows-native packaging for easy deployment",
        "Adaptive response using reinforcement learning",
        "Comprehensive dashboard with real-time visualization",
        "Self-Managing System (Autonomic Computing) - goes beyond simple monitoring"
    ])

    # Slide 12: Future Improvements
    add_content_slide("Future Improvements", [
        "Add user authentication and role-based access control",
        "Implement feedback loop for adaptive model learning (partially complete)",
        "Expand ML model ensemble with additional architectures",
        "Add mobile app for remote monitoring",
        "Integrate with SIEM systems",
        "Implement distributed deployment for enterprise networks",
        "Add threat intelligence feed integration"
    ])

    # Slide 13: Live Demo Overview
    add_content_slide("Live Demo Overview", [
        "1. Launch the application and access dashboard at http://127.0.0.1:5000",
        "2. Select network interface for packet capture",
        "3. Start live capture and observe real-time analysis",
        "4. View results table with model predictions and risk scores",
        "5. Demonstrate attack simulation capabilities",
        "6. Show prevention/healing actions (firewall blocking)",
        "7. Upload and analyze PCAP files",
        "8. Export results for further analysis"
    ])

    # Slide 14: Thank You
    add_title_slide(
        "Thank You",
        "Questions?\n\nCSCI480 Layered IDS/IPS Team"
    )

    # Save presentation
    output_path = r"c:\Users\pompk\Desktop\CSCI480\CSCI480_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
