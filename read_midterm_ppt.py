from pptx import Presentation

def read_pptx(file_path):
    prs = Presentation(file_path)
    
    print(f"Total slides: {len(prs.slides)}")
    print("=" * 80)
    
    for i, slide in enumerate(prs.slides):
        print(f"\nSLIDE {i + 1}")
        print("-" * 80)
        
        # Get slide title
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text}")
        
        # Get all text content
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape != slide.shapes.title:  # Skip title (already printed)
                    print(f"Content: {shape.text}")
        
        print()

if __name__ == "__main__":
    pptx_path = r"c:\Users\pompk\Downloads\_Midterm Alpha Prototype Review (3).pptx"
    read_pptx(pptx_path)
